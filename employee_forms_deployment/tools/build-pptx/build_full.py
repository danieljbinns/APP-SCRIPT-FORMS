# build_full.py - Generates PROJECT_SUMMARY_SLIDES.pptx from the full project summary.

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Colours ──────────────────────────────────────────────────────────────────
BG      = RGBColor(0x0c, 0x0c, 0x0c)
SURFACE = RGBColor(0x1a, 0x1a, 0x1a)
SURF2   = RGBColor(0x1e, 0x1e, 0x1e)
BORDER  = RGBColor(0x2e, 0x2e, 0x2e)
RED     = RGBColor(0xEB, 0x1C, 0x2D)
GREEN   = RGBColor(0x00, 0xC8, 0x53)
BLUE    = RGBColor(0x21, 0x96, 0xF3)
YELLOW  = RGBColor(0xFF, 0xAB, 0x00)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
TEXT    = RGBColor(0xEC, 0xEC, 0xEC)
MUTED   = RGBColor(0x90, 0x90, 0x90)

W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H
blank = prs.slide_layouts[6]


# ── Helpers ───────────────────────────────────────────────────────────────────

def set_bg(slide, color=BG):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def rect(slide, l, t, w, h, fill=None, line=None, lw=None):
    s = slide.shapes.add_shape(1, l, t, w, h)
    if fill:  s.fill.solid(); s.fill.fore_color.rgb = fill
    else:     s.fill.background()
    if line:  s.line.color.rgb = line;  s.line.width = lw or Pt(1)
    else:     s.line.fill.background()
    return s

def tb(slide, l, t, w, h, text, sz=14, bold=False, color=TEXT,
       align=PP_ALIGN.LEFT, italic=False, wrap=True):
    box = slide.shapes.add_textbox(l, t, w, h)
    box.word_wrap = wrap
    tf = box.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(sz)
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = 'Segoe UI'
    run.font.color.rgb = color
    return box

def header(slide, right=''):
    rect(slide, 0, 0, W, Inches(0.06), fill=RED)
    rect(slide, 0, Inches(0.06), W, Inches(0.55), fill=RGBColor(0x14, 0x14, 0x14))
    tb(slide, Inches(0.5), Inches(0.08), Inches(7), Inches(0.5),
       'TEAM Group  *  Employee Management  *  Project Summary', 13, bold=True, color=WHITE)
    if right:
        tb(slide, Inches(7), Inches(0.08), Inches(5.8), Inches(0.5),
           right, 12, color=MUTED, align=PP_ALIGN.RIGHT)

def slide_title(slide, title, sub=''):
    rect(slide, Inches(0.5), Inches(0.82), Inches(0.06), Inches(0.52), fill=RED)
    tb(slide, Inches(0.65), Inches(0.79), Inches(12.2), Inches(0.58),
       title, 26, bold=True, color=WHITE)
    if sub:
        tb(slide, Inches(0.65), Inches(1.38), Inches(12.2), Inches(0.36),
           sub, 13, italic=True, color=MUTED)

def tbl(slide, l, t, w, headers, rows, col_widths=None,
        hl_rows=None, total_rows=None):
    n_cols = len(headers)
    n_rows = len(rows) + 1
    row_h  = Inches(0.36)
    shape  = slide.shapes.add_table(n_rows, n_cols, l, t, w, row_h * n_rows)
    table  = shape.table

    if col_widths:
        for i, cw in enumerate(col_widths): table.columns[i].width = cw
    else:
        per = w // n_cols
        for i in range(n_cols): table.columns[i].width = per

    def cell(r, c, text, bg=SURFACE, fg=TEXT, bold=False, sz=11):
        cell_obj = table.cell(r, c)
        cell_obj.text = text
        p = cell_obj.text_frame.paragraphs[0]
        run = p.runs[0] if p.runs else p.add_run()
        run.text = text
        run.font.size = Pt(sz)
        run.font.bold = bold
        run.font.name = 'Segoe UI'
        run.font.color.rgb = fg
        cell_obj.fill.solid()
        cell_obj.fill.fore_color.rgb = bg

    for ci, h in enumerate(headers):
        cell(0, ci, h, bg=SURF2, fg=MUTED, bold=True, sz=10)

    for ri, row in enumerate(rows):
        is_hl    = hl_rows    and ri in hl_rows
        is_total = total_rows and ri in total_rows
        for ci, val in enumerate(row):
            if is_total:
                bg = RGBColor(0x28, 0x06, 0x08)
                fg = WHITE
                bold = True
            elif is_hl:
                bg = RGBColor(0x07, 0x20, 0x10)
                fg = GREEN if ci == 0 else TEXT
                bold = ci == 0
            else:
                bg = SURFACE if ri % 2 == 0 else RGBColor(0x16, 0x16, 0x16)
                fg = WHITE if ci == 0 else TEXT
                bold = ci == 0
            cell(ri + 1, ci, str(val), bg=bg, fg=fg, bold=bold, sz=11)

    return table


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Cover
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
set_bg(s)
header(s, 'April 2026  *  Internal Use Only')

rect(s, 0, Inches(0.61), W, Inches(0.04), fill=RED)

tb(s, Inches(1.5), Inches(1.2), Inches(10.33), Inches(0.3),
   'TEAM GROUP  *  INTERNAL TOOLING  *  APRIL 2026',
   11, color=MUTED, align=PP_ALIGN.CENTER)

tb(s, Inches(0.8), Inches(1.6), Inches(11.73), Inches(1.0),
   'Employee HR Workflow System', 44, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
tb(s, Inches(0.8), Inches(2.55), Inches(11.73), Inches(0.55),
   'Project Summary', 32, bold=True, color=RED, align=PP_ALIGN.CENTER)

rect(s, Inches(5.67), Inches(3.25), Inches(2.0), Inches(0.04), fill=RED)

tb(s, Inches(1.5), Inches(3.4), Inches(10.33), Inches(1.6),
   'From a blank repository to a production multi-step HR workflow engine spanning\n'
   'three Google Apps Script environments — built entirely inside Google Workspace.\n'
   'Zero external servers.  Zero paid SaaS.  Zero database.',
   16, color=TEXT, align=PP_ALIGN.CENTER)

for i, pill in enumerate(['331 tracked changes', '~174 hours', '48 active dev days',
                           '4 workflow types', '105 staging workflows']):
    pw = Inches(2.35)
    gap = Inches(0.1)
    total_w = 5 * pw + 4 * gap
    lx = (W - total_w) / 2 + i * (pw + gap)
    rect(s, lx, Inches(5.2), pw, Inches(0.44), fill=SURF2)
    tb(s, lx, Inches(5.22), pw, Inches(0.4),
       pill, 12, bold=True, color=MUTED, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — At a Glance
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
set_bg(s)
header(s, 'At a Glance')
slide_title(s, 'At a Glance', '5 months  *  solo  *  part-time  *  Dec 2025 – Apr 2026')

KPIs = [
    ('331',    YELLOW, 'Tracked Changes',   '128 git commits + 203 clasp saves'),
    ('273K',   BLUE,   'Lines Added',       '~248K net written'),
    ('232',    BLUE,   'Source Files',      '3 GCP projects, 3 GAS environments'),
    ('48',     GREEN,  'Active Dev Days',   'Coding + QA — all dev work'),
    ('~174h',  GREEN,  'Human Hours',       '~133h dev + ~41h QA'),
    ('4',      RED,    'Workflow Types',    'New Hire, TERM/EOE, Equipment, Status Change'),
    ('105',    BLUE,   'Staging Workflows', 'All developer QA submissions'),
    ('~305',   YELLOW, 'Sheet Writes',      'Confirmed end-to-end executions'),
]

card_w = Inches(2.95)
card_h = Inches(1.55)
gap_x  = Inches(0.2)
gap_y  = Inches(0.18)
start_l = Inches(0.62)
start_t = Inches(1.88)

for i, (num, color, label, sub) in enumerate(KPIs):
    col = i % 4
    row = i // 4
    l = start_l + col * (card_w + gap_x)
    t = start_t + row * (card_h + gap_y)
    rect(s, l, t, card_w, card_h, fill=SURFACE)
    tb(s, l, t + Inches(0.18), card_w, Inches(0.65),
       num, 46, bold=True, color=color, align=PP_ALIGN.CENTER)
    tb(s, l, t + Inches(0.85), card_w, Inches(0.28),
       label.upper(), 10, bold=True, color=MUTED, align=PP_ALIGN.CENTER)
    tb(s, l, t + Inches(1.14), card_w, Inches(0.3),
       sub, 10, color=MUTED, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — What Was Built
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
set_bg(s)
header(s, 'What Was Built')
slide_title(s, 'What Was Built',
            'Replaced emailed spreadsheets with a gated, multi-step request pipeline.')

# Workflow type pills
tb(s, Inches(0.62), Inches(1.88), Inches(3), Inches(0.26),
   'WORKFLOW TYPES', 10, bold=True, color=RED)

TYPES = ['New Hire', 'Termination / End of Employment (TERM + EOE)',
         'Equipment Request', 'Status / Position Change']
tw = Inches(2.9)
for i, wt in enumerate(TYPES):
    lx = Inches(0.62) + i * (tw + Inches(0.18))
    rect(s, lx, Inches(2.14), tw, Inches(0.42), fill=SURF2)
    rect(s, lx, Inches(2.14), Inches(0.05), Inches(0.42), fill=RED)
    tb(s, lx + Inches(0.1), Inches(2.16), tw - Inches(0.15), Inches(0.38),
       wt, 12, bold=True, color=WHITE)

tbl(s,
    l=Inches(0.5), t=Inches(2.72), w=Inches(12.33),
    headers=['Capability', 'Detail'],
    rows=[
        ['Sequential gating',        'Each step unlocks only after the prior completes — no skipping, no out-of-order, enforced by the system'],
        ['Role-based dashboard',      'HR, IT, Safety, and Managers each see their own view via Google Group membership — no new accounts'],
        ['Automatic task routing',    'Credit Card, Fleetio, Jonas, SiteDocs, BOSS, ADP, 30-60-90 — each specialist team notified and routed automatically'],
        ['Unified Action Items',      '6 legacy specialist sheets retired; replaced with a single normalized task table with full lifecycle tracking'],
        ['Form edit mode',            'HR Verification and IT Confirmation can re-open and modify data — full before/after change-audit email dispatched'],
        ['Full audit trail',          'Every action timestamped, every edit notified. Complete accountability at every step from request to close'],
        ['Environment-aware config',  'Same codebase deploys to prod, staging, and dev via Script Properties — no hardcoded IDs, safe to migrate'],
    ],
    col_widths=[Inches(2.5), Inches(9.83)],
)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — Change History & Scale
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
set_bg(s)
header(s, 'Change History & Scale')
slide_title(s, 'Change History & Scale',
            '331 tracked changes across 48 active dev days — Dec 2025 through Apr 2026.')

# Left: key metrics table
tbl(s,
    l=Inches(0.5), t=Inches(1.88), w=Inches(4.8),
    headers=['Metric', 'Value'],
    rows=[
        ['Total tracked changes',     '331'],
        ['Active dev days',           '48'],
        ['Total lines added',         '273,598'],
        ['Net lines written',         '~248,425'],
        ['Unique file paths touched', '1,426'],
        ['Active span',               '~147 days (5 months)'],
        ['First active day',          '2025-12-04'],
        ['Last active day',           '2026-04-29'],
    ],
    hl_rows=[0, 1],
    col_widths=[Inches(2.6), Inches(2.2)],
)

# Right: commits by month
tbl(s,
    l=Inches(5.5), t=Inches(1.88), w=Inches(7.33),
    headers=['Month', 'Commits', 'Phase'],
    rows=[
        ['Dec 2025', '20',  'Foundations, architecture, initial 3-form system'],
        ['Jan 2026', '26',  'V2 architecture, full system init, prod deployment'],
        ['Feb 2026', '17',  'Environment-aware config, staging + dev stood up'],
        ['Mar 2026', '91',  'Equipment requests, dashboard optimization, staging marathon'],
        ['Apr 2026', '177', 'Action Items refactor + all new workflow types + full QA sprint'],
        ['Total',    '331', '48 active dev days'],
    ],
    hl_rows=[4],
    total_rows=[5],
    col_widths=[Inches(1.3), Inches(1.0), Inches(5.03)],
)

# April callout
rect(s, Inches(5.5), Inches(5.62), Inches(7.33), Inches(0.88),
     fill=RGBColor(0x20, 0x08, 0x08))
tb(s, Inches(5.65), Inches(5.68), Inches(7.1), Inches(0.78),
   'April sprint: 177 of 331 commits (53% of the project) in 10 days. Peak day Apr 21: '
   '32 commits — Status Change, EOE, form edit mode, and ChangeNotify all in one session.',
   12, color=TEXT)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — Codebase by Environment
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
set_bg(s)
header(s, 'Codebase by Environment')
slide_title(s, 'Codebase by Environment',
            '3 isolated GCP projects — prod, staging, dev — same codebase, environment-aware config.')

ENVS = [
    ('Production',  'employee_management_v2/',         'Live', RED,
     [('HTML files', '16'), ('JS files', '0 (HTML-embedded)'), ('Lines of code', '~7,393'), ('State', 'Live on prod')]),
    ('Staging',     'employee_management_v2_staging/',  'QA',  BLUE,
     [('JS files', '21'), ('HTML files', '29'), ('Service manifests', '4'), ('Lines of code', '~18,639')]),
    ('Dev',         'employee_management_v2_dev/',      'Next', YELLOW,
     [('JS files', '23'), ('HTML files', '31'), ('Service manifests', '4'), ('Lines of code', '~19,355')]),
]

env_w = Inches(3.95)
for i, (name, path, badge, color, rows) in enumerate(ENVS):
    l = Inches(0.5) + i * (env_w + Inches(0.24))
    rect(s, l, Inches(1.88), env_w, Inches(2.45), fill=SURFACE)
    rect(s, l, Inches(1.88), env_w, Inches(0.05), fill=color)
    tb(s, l + Inches(0.18), Inches(1.96), env_w - Inches(0.3), Inches(0.32),
       f'{name}  [{badge}]', 15, bold=True, color=WHITE)
    tb(s, l + Inches(0.18), Inches(2.3), env_w - Inches(0.3), Inches(0.22),
       path, 10, color=MUTED)
    for j, (k, v) in enumerate(rows):
        ty = Inches(2.56) + j * Inches(0.42)
        tb(s, l + Inches(0.18), ty, env_w * 0.5, Inches(0.38), k, 11, color=MUTED)
        tb(s, l + env_w * 0.5, ty, env_w * 0.5 - Inches(0.2), Inches(0.38),
           v, 11, bold=True, color=WHITE, align=PP_ALIGN.RIGHT)

# Key files table
tbl(s,
    l=Inches(0.5), t=Inches(4.52), w=Inches(12.33),
    headers=['File', '~Lines', 'Purpose'],
    rows=[
        ['Dashboard.html',               '2,400', 'Main operator dashboard — role-gated, multi-view, live status + process flowchart'],
        ['WorkflowEngine.js',            '600',   'Sequential gating, step transitions, state machine — powers all 4 workflow types'],
        ['EmailUtils.js + ChangeNotify.js', '~930', 'Email templates, subject builder, change-diff audit emails, conditional CC routing'],
        ['ActionItemsManager.js',        '450',   'Unified task CRUD, lifecycle (Pending→Closed), specialist routing'],
        ['InitialRequest.html',          '900',   'New hire request form — dual-list pickers, BOSS, Jonas cost sheets'],
    ],
    col_widths=[Inches(2.8), Inches(0.9), Inches(8.63)],
)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — Human Hours
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
set_bg(s)
header(s, 'Human Hours')
slide_title(s, 'Human Hours — ~174 h',
            '~133 h development  *  ~41 h QA  *  48 active dev days  *  solo, part-time.')

# Monthly summary table
tbl(s,
    l=Inches(0.5), t=Inches(1.88), w=Inches(7.5),
    headers=['Month', 'Tracked Changes', 'Total Hours', 'Dev Hours', 'QA Hours'],
    rows=[
        ['Dec 2025', '20',  '~15 h',  '~15 h', '—'],
        ['Jan 2026', '26',  '~28 h',  '~24 h', '~4 h'],
        ['Feb 2026', '17',  '~10 h',  '~7 h',  '~3 h'],
        ['Mar 2026', '91',  '~46 h',  '~38 h', '~8 h'],
        ['Apr 2026', '177', '~75 h',  '~49 h', '~26 h'],
        ['Total — 48 dev days', '331', '~174 h', '~133 h', '~41 h'],
    ],
    hl_rows=[4],
    total_rows=[5],
    col_widths=[Inches(2.0), Inches(1.35), Inches(1.35), Inches(1.4), Inches(1.4)],
)

# April velocity table
tbl(s,
    l=Inches(8.1), t=Inches(1.88), w=Inches(4.73),
    headers=['Date', 'Commits', 'Notable Work'],
    rows=[
        ['Apr 14', '2',  'Phase 8 QA — regressions R1–R5'],
        ['Apr 16', '26', 'Central Purchasing, buildEmailSubject, Groups A–F'],
        ['Apr 17', '27', 'Staging hotfixes — empType, dates, HR verification'],
        ['Apr 20', '12', 'Action Items refactor phases 1–6'],
        ['Apr 21', '32', 'STATUS CHANGE + EOE + form edit + ChangeNotify'],
        ['Apr 22', '16', 'Equipment form, Jonas/Central Purchasing fix'],
        ['Apr 23', '12', 'Dual-list picker, BOSS gatekeeper, full QA pass'],
        ['Apr 27', '20', 'Major dev sprint [dev]'],
        ['Apr 28', '22', 'Email subject overflow fix, EQUIP_REQ dashboard'],
        ['Apr 29', '3',  'Final dev push [dev]'],
    ],
    hl_rows=[4],
    col_widths=[Inches(0.8), Inches(0.8), Inches(3.13)],
)

# Cost avoidance note
rect(s, Inches(0.5), Inches(5.55), Inches(7.5), Inches(1.32),
     fill=RGBColor(0x04, 0x18, 0x0c))
tb(s, Inches(0.65), Inches(5.62), Inches(7.25), Inches(1.2),
   'Cost avoidance: outsourcing an equivalent system = $65K–$115K one-time. '
   'SaaS workflow platforms = $10–$25/user/mo ($180K–$450K+/yr at scale). '
   'Basic HR tools = $6–$9/user/mo ($108K–$162K/yr, data collection only). '
   'This runs on Google Workspace TEAM Group already pays for — zero incremental cost.',
   12, color=TEXT)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — QAT & Testing Evidence
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
set_bg(s)
header(s, 'QAT & Testing Evidence')
slide_title(s, 'QAT & Testing Evidence',
            '6 deliberate QA sessions  *  105 staging workflows  *  ~305 confirmed sheet writes.')

# Testing phases callout
rect(s, Inches(0.5), Inches(1.88), Inches(12.33), Inches(0.84),
     fill=SURF2)
tb(s, Inches(0.65), Inches(1.93), Inches(12.0), Inches(0.76),
   'Phase 1 — Prod-only (Jan 23 – Feb 24): No staging yet. ~72 NEW_EMP workflows on prod.   '
   'Phase 2 — Staging primary (Feb 24 – Apr 20): Prod cleaned; staging became active QA surface.   '
   'Phase 3 — Dev for refactor (Apr 21–29): Dev seeded with prod data; all 4 types on new schema.',
   12, color=TEXT)

# QA sessions table
tbl(s,
    l=Inches(0.5), t=Inches(2.86), w=Inches(12.33),
    headers=['Session', 'Date', 'Duration', 'Clasp Saves', 'Scope / Output'],
    rows=[
        ['New Hire prod rollout',    'Jan 23, 2026',  '~3.5 h', '13 prod',     'Full NEW_EMP walkthrough — produced ~72 test records later migrated to staging'],
        ['Staging standup + cutover','Feb 24, 2026',  '~3 h',   '10 staging',  'Staging & dev created; prod test data moved; prod cleaned for real HR use only'],
        ['Termination / EOE sprint', 'Mar 6–9, 2026', '~7.5 h', '37 staging',  'Full TERM + EOE workflow — gating, Action Items, approval routing, email edge cases'],
        ['Action Items refactor QA', 'Apr 21, 2026',  '~6.7 h', '15 dev',      'Unified schema — all 4 workflow types on new data model (14:45–21:29 UTC)'],
        ['Final validation I',       'Apr 27, 2026',  '~7.5 h', '20 dev',      'Cross-env check; email edge cases; status change pipeline'],
        ['Final validation II',      'Apr 28, 2026',  '~4.8 h', '18 dev',      'IT confirmation gate; Equipment Request; pre-launch final checks'],
    ],
    hl_rows=[2, 3],
    col_widths=[Inches(2.2), Inches(1.6), Inches(1.0), Inches(1.15), Inches(6.38)],
)

# Staging test record
tbl(s,
    l=Inches(0.5), t=Inches(5.7), w=Inches(7.5),
    headers=['Sheet', 'Records', 'Notes'],
    rows=[
        ['Workflows (total)',  '105', '72 new hire + 30 TERM/EOE + 3 position change'],
        ['Action Items',       '117', 'Generated from TERM/EOE/CHANGE workflows'],
        ['Total sheet writes', '~305','Each = one successful end-to-end execution'],
    ],
    total_rows=[2],
    col_widths=[Inches(2.4), Inches(1.0), Inches(4.1)],
)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — Development Timeline
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
set_bg(s)
header(s, 'Development Timeline')
slide_title(s, 'Development Timeline',
            'Dec 2025 → Apr 2026 — from blank repo to production system.')

MILESTONES = [
    ('Dec 2025',  RED,    '20 commits',  [
        ('Dec 04', 'Initial commit — REQUEST_FORMS demo system (complete prototype)'),
        ('Dec 09', 'MAJOR REORGANIZATION — clean modular structure, clear module boundaries'),
        ('Dec 17', 'Deploy all 9 forms to shared drive, clasp pipeline, Styles library'),
    ]),
    ('Jan 2026',  BLUE,   '26 commits',  [
        ('Jan 22', 'Initialize comprehensive employee management system (V2) — full module structure'),
        ('Jan 23', 'First prod deployment & new hire testing — 13 iterative pushes [prod]'),
        ('Jan 27', 'Safety Notification & Config Cleanup / Safety + HR Email Fixes [prod]'),
    ]),
    ('Feb 2026',  YELLOW, '17 commits',  [
        ('Feb 20', 'Environment-aware config via Script Properties — no more hardcoded IDs'),
        ('Feb 24', 'Staging + dev environments stood up from scratch — 15 commits [staging+dev]'),
    ]),
    ('Mar 2026',  GREEN,  '91 commits',  [
        ('Mar 06', 'Major staging feature build — 27 commits, specialist forms, dashboard [staging]'),
        ('Mar 25', 'Equipment requests workflow, dashboard performance optimization [staging+prod]'),
    ]),
    ('Apr 2026',  RED,    '177 commits', [
        ('Apr 20', 'ACTION ITEMS REFACTOR — phases 1–6 (6 legacy sheets → unified task system)'),
        ('Apr 21', 'Status Change + EOE workflow + Form edit mode + ChangeNotify (32 commits — peak day)'),
        ('Apr 23', 'Dual-list picker, BOSS gatekeeper, full QA pass across all forms'),
        ('Apr 28', 'Email subject overflow fix, EQUIP_REQ dashboard, pre-launch checks'),
    ]),
]

col_w = Inches(2.4)
gap   = Inches(0.16)
t_start = Inches(1.88)

for i, (month, color, commits, items) in enumerate(MILESTONES):
    l = Inches(0.5) + i * (col_w + gap)
    # Month header
    rect(s, l, t_start, col_w, Inches(0.5), fill=RGBColor(0x20, 0x20, 0x20))
    rect(s, l, t_start, col_w, Inches(0.05), fill=color)
    tb(s, l + Inches(0.1), t_start + Inches(0.06), col_w - Inches(0.15), Inches(0.2),
       month, 13, bold=True, color=WHITE)
    tb(s, l + Inches(0.1), t_start + Inches(0.27), col_w - Inches(0.15), Inches(0.18),
       commits, 10, color=MUTED)
    # Items
    for j, (date, desc) in enumerate(items):
        ty = t_start + Inches(0.58) + j * Inches(1.1)
        if ty + Inches(1.0) > Inches(7.2): break
        rect(s, l, ty, col_w, Inches(1.0), fill=SURF2)
        rect(s, l, ty, Inches(0.04), Inches(1.0), fill=color)
        tb(s, l + Inches(0.1), ty + Inches(0.06), col_w - Inches(0.15), Inches(0.2),
           date, 10, color=MUTED)
        tb(s, l + Inches(0.1), ty + Inches(0.26), col_w - Inches(0.15), Inches(0.68),
           desc, 11, color=TEXT)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — Architecture Highlights
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
set_bg(s)
header(s, 'Architecture Highlights')
slide_title(s, 'Architecture Highlights',
            'Enterprise workflow capabilities — inside a Google Sheet.')

ARCHES = [
    ('Action Items Refactor',
     'The single largest architectural decision. Six independent specialist result sheets retired '
     'and replaced with one normalized Action Items tab. Every task = one row. Lifecycle uniform. '
     'Form data stored as JSON blob — schema changes need no migration. '
     'Migration ran as a local Node.js script (dry-run / execute) — zero changes to live GAS code.'),
    ('Email Engine (~1,200 lines)',
     'EmailUtils.js + ChangeNotify.js: dark-mode inline-styled HTML templates, per-workflow subject '
     'builder with 250-char hard cap (hit in prod on Status Change — long site names + classification '
     'strings pushed subject to ~300 chars), change-audit diff tables (field / was / now), '
     'conditional CC routing for Safety and ID Setup on relevant changes.'),
    ('Sequential Gating Engine',
     'Each workflow defines its steps as a config object. Engine evaluates state, identifies next '
     'pending step, generates gated form URL, refuses to advance if prior steps are incomplete. '
     'Same engine powers all 4 workflow types with zero duplication of gating logic.'),
]

card_w = Inches(3.9)
gap_x  = Inches(0.26)
for i, (title, body) in enumerate(ARCHES):
    l = Inches(0.5) + i * (card_w + gap_x)
    rect(s, l, Inches(1.88), card_w, Inches(2.9), fill=SURFACE)
    rect(s, l, Inches(1.88), card_w, Inches(0.05), fill=RED)
    tb(s, l + Inches(0.18), Inches(1.96), card_w - Inches(0.3), Inches(0.36),
       title, 14, bold=True, color=WHITE)
    tb(s, l + Inches(0.18), Inches(2.36), card_w - Inches(0.3), Inches(2.3),
       body, 11, color=RGBColor(0xd0, 0xd0, 0xd0))

# "What makes this unusual" table
tbl(s,
    l=Inches(0.5), t=Inches(4.95), w=Inches(12.33),
    headers=['Typically a GAS script does...', 'This project does...'],
    rows=[
        ['A form submission hook',         'A full multi-step, multi-role workflow state machine across 4 workflow types'],
        ['A spreadsheet helper formula',   'A normalized relational data model across 10+ sheets with consistent IDs and joins'],
        ['A simple HTML sidebar',          'A live web app with session state, role checks, and client-server RPC via google.script.run'],
        ['One GCP project, one clasp ID',  'Three isolated GCP projects, environment-aware config, local migration tooling for safe upgrades'],
    ],
    col_widths=[Inches(3.8), Inches(8.53)],
)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — Key Technical Decisions
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
set_bg(s)
header(s, 'Key Technical Decisions')
slide_title(s, 'Key Technical Decisions',
            'Every decision here solved a real constraint — version-controlled, maintainable, safe to hand off.')

tbl(s,
    l=Inches(0.5), t=Inches(1.88), w=Inches(12.33),
    headers=['Decision', 'Rationale'],
    rows=[
        ['.js over .gs everywhere',       'Editors treat .js as JavaScript properly; cleaner clasp push/pull diffs'],
        ['JSON blobs for form data',       'Schema-flexible — specialist forms can evolve without a sheet column migration'],
        ['Script Properties for config',   'Environment switching without code changes; no hardcoded spreadsheet or script IDs'],
        ['Local migration tooling',        'Prod GAS untouched during data model migrations; safe dry-run workflow before execute'],
        ['Inline styles in emails',        'Google Workspace strips <style> tags in some email clients; inline-only is reliable'],
        ['250-char subject cap',           'GAS throws Argument too large: subject above ~250 chars — hit in prod on Status Change'],
        ['EPPlus for Excel migration',     'Import-Excel fails on duplicate column headers; EPPlus positional copy bypasses this'],
        ['Separate GCP projects per env',  'Prod OAuth scopes and service account permissions isolated from staging/dev experiments'],
    ],
    col_widths=[Inches(3.0), Inches(9.33)],
)

# Bottom line
rect(s, Inches(0.5), Inches(5.82), Inches(12.33), Inches(1.05),
     fill=RGBColor(0x04, 0x18, 0x0c))
tb(s, Inches(0.65), Inches(5.88), Inches(12.0), Inches(0.92),
   'TEAM Group HR onboarding and offboarding moved from emailed spreadsheets to a fully tracked, '
   'accountable, automated workflow — built on tools the organization already uses, with zero new software, '
   'zero recurring licensing cost, and no external contractors. '
   'Equivalent capability: $65K–$115K to outsource  |  $180K–$450K+/yr to license.',
   13, color=TEXT)


# ── Save ─────────────────────────────────────────────────────────────────────
out = r"P:\Repos\github\danieljbinns\APP SCRIPT FORMS\employee_forms_deployment\PROJECT_SUMMARY_SLIDES.pptx"
prs.save(out)
print(f"Saved -> {out}")
print(f"Slides: {len(prs.slides)}")
