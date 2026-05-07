"""
build.py — Generates PROJECT_SUMMARY_EXEC_SLIDES.pptx
Matches the HTML slide deck: dark theme, TEAM Group branding, 10 slides.
Run: python build.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from pptx.oxml.ns import qn
from pptx.dml.color import RGBColor
import copy
from lxml import etree

# ── Colours ──────────────────────────────────────────────────────────────────
BG      = RGBColor(0x0c, 0x0c, 0x0c)
SURFACE = RGBColor(0x1a, 0x1a, 0x1a)
BORDER  = RGBColor(0x2e, 0x2e, 0x2e)
RED     = RGBColor(0xEB, 0x1C, 0x2D)
GREEN   = RGBColor(0x00, 0xC8, 0x53)
BLUE    = RGBColor(0x21, 0x96, 0xF3)
YELLOW  = RGBColor(0xFF, 0xAB, 0x00)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
TEXT    = RGBColor(0xEC, 0xEC, 0xEC)
MUTED   = RGBColor(0x90, 0x90, 0x90)
DARK_RED= RGBColor(0x2e, 0x00, 0x04)

# ── Slide dimensions (widescreen 16:9) ───────────────────────────────────────
W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

blank_layout = prs.slide_layouts[6]   # truly blank


# ── Helpers ───────────────────────────────────────────────────────────────────

def set_bg(slide, color):
    """Fill slide background with a solid colour."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_rect(slide, l, t, w, h, fill_color=None, line_color=None, line_width=None):
    shape = slide.shapes.add_shape(1, l, t, w, h)  # MSO_SHAPE_TYPE.RECTANGLE = 1
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        if line_width:
            shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape

def add_textbox(slide, l, t, w, h, text, font_size=18, bold=False, color=None,
                align=PP_ALIGN.LEFT, italic=False, wrap=True, font_name='Segoe UI'):
    txb = slide.shapes.add_textbox(l, t, w, h)
    txb.word_wrap = wrap
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = font_name
    if color:
        run.font.color.rgb = color
    return txb

def add_para(tf, text, font_size=14, bold=False, color=None, italic=False,
             align=PP_ALIGN.LEFT, space_before=0):
    p = tf.add_paragraph()
    p.alignment = align
    if space_before:
        p.space_before = Pt(space_before)
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = 'Segoe UI'
    if color:
        run.font.color.rgb = color
    return p

def add_header(slide, slide_name=''):
    """Red-bordered top bar + branding."""
    # Red accent line
    add_rect(slide, 0, 0, W, Inches(0.06), fill_color=RED)
    # Brand bar
    add_rect(slide, 0, Inches(0.06), W, Inches(0.55), fill_color=RGBColor(0x14,0x14,0x14))
    add_textbox(slide, Inches(0.5), Inches(0.08), Inches(6), Inches(0.5),
                'TEAM Group  ·  Employee Management', 13, bold=True, color=WHITE)
    right_text = slide_name if slide_name else 'April 2026  ·  Internal Use Only'
    add_textbox(slide, Inches(7.5), Inches(0.08), Inches(5.3), Inches(0.5),
                right_text, 12, color=MUTED, align=PP_ALIGN.RIGHT)

def add_slide_title(slide, title, subtitle=''):
    """Red-left-border section heading + italic subtitle."""
    # Red bar (left border effect)
    add_rect(slide, Inches(0.5), Inches(0.85), Inches(0.06), Inches(0.5), fill_color=RED)
    add_textbox(slide, Inches(0.65), Inches(0.82), Inches(12.2), Inches(0.55),
                title, 28, bold=True, color=WHITE)
    if subtitle:
        add_textbox(slide, Inches(0.65), Inches(1.38), Inches(12.2), Inches(0.38),
                    subtitle, 14, italic=True, color=MUTED)

def add_table(slide, l, t, w, headers, rows, col_widths=None,
              hl_rows=None, total_row=False):
    """
    Add a styled table.  hl_rows = list of row indices (0-based, excluding header) to highlight green.
    total_row = True means last row is bold red-ish.
    """
    n_cols = len(headers)
    n_rows = len(rows) + 1  # +1 for header
    row_h  = Inches(0.38)
    tbl_h  = row_h * n_rows

    tbl = slide.shapes.add_table(n_rows, n_cols, l, t, w, tbl_h).table

    # Column widths
    if col_widths:
        for i, cw in enumerate(col_widths):
            tbl.columns[i].width = cw
    else:
        per = w // n_cols
        for i in range(n_cols):
            tbl.columns[i].width = per

    def style_cell(cell, text, bg=SURFACE, fg=TEXT, bold=False, sz=11):
        cell.text = text
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        run = p.runs[0] if p.runs else p.add_run()
        run.text = text
        run.font.size = Pt(sz)
        run.font.bold = bold
        run.font.name = 'Segoe UI'
        run.font.color.rgb = fg
        cell.fill.solid()
        cell.fill.fore_color.rgb = bg

    # Header row
    for ci, h in enumerate(headers):
        style_cell(tbl.cell(0, ci), h,
                   bg=RGBColor(0x20, 0x20, 0x20), fg=MUTED, bold=True, sz=10)

    # Data rows
    for ri, row in enumerate(rows):
        is_hl    = hl_rows and ri in hl_rows
        is_total = total_row and ri == len(rows) - 1
        for ci, val in enumerate(row):
            if is_total:
                bg = RGBColor(0x28, 0x06, 0x08)
                fg = WHITE if ci == 0 else TEXT
                bold = True
            elif is_hl:
                bg = RGBColor(0x07, 0x20, 0x10)
                fg = GREEN if ci == 0 else TEXT
                bold = ci == 0
            else:
                bg = SURFACE if ri % 2 == 0 else RGBColor(0x16, 0x16, 0x16)
                fg = WHITE if ci == 0 else TEXT
                bold = ci == 0
            style_cell(tbl.cell(ri + 1, ci), str(val), bg=bg, fg=fg, bold=bold, sz=11)

    return tbl


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Cover
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
set_bg(s, BG)
add_header(s)

# Centre content
add_textbox(s, Inches(1.5), Inches(1.4), Inches(10.33), Inches(0.35),
            'TEAM GROUP  ·  INTERNAL TOOLING  ·  APRIL 2026',
            11, color=MUTED, align=PP_ALIGN.CENTER)

# Title
add_textbox(s, Inches(1.0), Inches(1.85), Inches(11.33), Inches(1.1),
            'Employee HR Workflow System', 44, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_textbox(s, Inches(1.0), Inches(2.85), Inches(11.33), Inches(0.65),
            'Executive Summary', 36, bold=True, color=RED, align=PP_ALIGN.CENTER)

# Red divider
add_rect(s, Inches(5.67), Inches(3.6), Inches(2.0), Inches(0.04), fill_color=RED)

# Tagline
add_textbox(s, Inches(1.5), Inches(3.75), Inches(10.33), Inches(1.3),
            'A production-grade, multi-step HR workflow engine built entirely inside Google Workspace —\n'
            'zero external servers, zero paid SaaS, zero database.\n'
            'Four workflow types. Full role-based access. Complete audit trails.',
            16, color=TEXT, align=PP_ALIGN.CENTER)

# CTA
add_textbox(s, Inches(1.5), Inches(5.2), Inches(10.33), Inches(0.4),
            'Use arrow keys or click to navigate', 12, italic=True, color=MUTED, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — At a Glance
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
set_bg(s, BG)
add_header(s, 'At a Glance')
add_slide_title(s, 'At a Glance', 'Key numbers from a 5-month solo build')

KPIs = [
    ('~174',  GREEN,  'Human Hours',         '48 active days · Dec 2025–Apr 2026'),
    ('$0',    BLUE,   'Recurring Cost',       'vs. $180K–$450K+/yr equivalent platforms'),
    ('4',     RED,    'Workflow Types',       'Live in production'),
    ('331',   YELLOW, 'Tracked Changes',     '128 git commits · 203 clasp saves'),
]

cols = [0, 1]
rows_idx = [0, 1]
cell_w = Inches(5.8)
cell_h = Inches(2.2)
gap    = Inches(0.26)
start_l = Inches(0.62)
start_t = Inches(1.85)

for i, (num, color, label, sub) in enumerate(KPIs):
    col = i % 2
    row = i // 2
    l = start_l + col * (cell_w + gap)
    t = start_t + row * (cell_h + gap)
    # Card bg
    add_rect(s, l, t, cell_w, cell_h, fill_color=SURFACE)
    # Big number
    add_textbox(s, l, t + Inches(0.35), cell_w, Inches(0.9),
                num, 64, bold=True, color=color, align=PP_ALIGN.CENTER)
    # Label
    add_textbox(s, l, t + Inches(1.3), cell_w, Inches(0.35),
                label.upper(), 11, bold=True, color=MUTED, align=PP_ALIGN.CENTER)
    # Sub
    add_textbox(s, l, t + Inches(1.65), cell_w, Inches(0.4),
                sub, 12, color=MUTED, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — What Was Delivered
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
set_bg(s, BG)
add_header(s, 'What Was Delivered')
add_slide_title(s, 'What Was Delivered',
                'Six capabilities. One system. Zero new software. Zero added cost.')

FEATURES = [
    ('4 HR Workflow Types',
     'New Hire, Termination / End of Employment, Equipment Request, and '
     'Position / Site Change — each fully automated from first request through to completion.'),
    ('Multi-Step Approval Pipeline',
     'Every request moves through defined steps in order. No step can be skipped or '
     'completed twice — the system enforces the process, not individual discipline.'),
    ('Role-Based Dashboard',
     'HR, IT, Safety, and Managers each see only their own queue — controlled through '
     'existing Google Groups with no new accounts or permissions to manage.'),
    ('Automatic Task Routing & Departmental Notifications',
     'IT credentials, fleet access, credit cards, purchasing accounts, safety onboarding — '
     'the right team is notified and assigned automatically. No manual forwarding, no missed handoffs.'),
    ('Full Audit Trail & Change Tracking',
     'Every action is timestamped and logged. Any edit notifies all affected parties '
     'with a full before/after summary — complete accountability at every step.'),
    ('Everything Stays in Google Workspace',
     'All forms, workflow data, and history live in your existing Drive and Sheets — '
     'no vendor access to your HR data, no new platforms to learn, no integration contracts.'),
]

card_w = Inches(4.04)
card_h = Inches(1.9)
gap_x  = Inches(0.24)
gap_y  = Inches(0.18)
start_l = Inches(0.62)
start_t = Inches(1.85)

for i, (title, body) in enumerate(FEATURES):
    col = i % 3
    row = i // 3
    l = start_l + col * (card_w + gap_x)
    t = start_t + row * (card_h + gap_y)
    # Red top accent bar
    add_rect(s, l, t, card_w, Inches(0.05), fill_color=RED)
    # Card bg
    add_rect(s, l, t + Inches(0.05), card_w, card_h - Inches(0.05), fill_color=RGBColor(0x1e,0x1e,0x1e))
    # Card title
    add_textbox(s, l + Inches(0.18), t + Inches(0.12), card_w - Inches(0.3), Inches(0.38),
                title, 13, bold=True, color=WHITE)
    # Card body
    add_textbox(s, l + Inches(0.18), t + Inches(0.52), card_w - Inches(0.3), Inches(1.25),
                body, 11, color=RGBColor(0xd0, 0xd0, 0xd0))


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — The Real Win
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
set_bg(s, BG)
add_header(s, 'The Real Win')
add_slide_title(s, 'The Real Win',
                "The biggest win isn't the cost — it's workflows that actually get done and nothing that falls through.")

BEFORE = [
    'Excel or Word forms emailed between departments',
    'No visibility on who acted, or when',
    'Items lost in inboxes — no accountability',
    'Specialists informed by manual forwarding, if at all',
    'No audit trail for HR, compliance, or disputes',
    'Follow-ups relied entirely on individual memory',
]
AFTER = [
    'Single tracked request per employee, start to finish',
    'Every step timestamped — who acted and when',
    "Nothing can be skipped — the system enforces completion",
    "Specialists auto-notified the moment they're needed",
    'Full change history and audit trail on every record',
    'Zero manual follow-up for routine HR workflows',
]

col_w = Inches(6.0)
col_h = Inches(4.2)
t     = Inches(1.85)

# BEFORE column
add_rect(s, Inches(0.5), t, col_w, col_h, fill_color=RGBColor(0x14,0x14,0x14))
add_textbox(s, Inches(0.7), t + Inches(0.15), Inches(2), Inches(0.3),
            'THE OLD PROCESS', 10, bold=True, color=MUTED)
for i, item in enumerate(BEFORE):
    add_textbox(s, Inches(0.7), t + Inches(0.5) + i * Inches(0.58),
                col_w - Inches(0.4), Inches(0.52),
                '✕  ' + item, 12, color=TEXT)

# AFTER column
add_rect(s, Inches(6.7), t, col_w, col_h, fill_color=RGBColor(0x04, 0x18, 0x0c))
add_textbox(s, Inches(6.9), t + Inches(0.15), Inches(2), Inches(0.3),
            'NOW', 10, bold=True, color=GREEN)
for i, item in enumerate(AFTER):
    add_textbox(s, Inches(6.9), t + Inches(0.5) + i * Inches(0.58),
                col_w - Inches(0.4), Inches(0.52),
                '✓  ' + item, 12, color=TEXT)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — Value & Cost Avoidance
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
set_bg(s, BG)
add_header(s, 'Value & Cost Avoidance')
add_slide_title(s, 'Value & Cost Avoidance',
                'Cost avoidance — not hypothetical savings. Real numbers from real platform pricing.')

# Left bar — SaaS cost
bar_t = Inches(1.85)
bar_h = Inches(1.7)
bar_w = Inches(6.1)

add_rect(s, Inches(0.5), bar_t, bar_w, bar_h, fill_color=RGBColor(0x20, 0x18, 0x04))
add_textbox(s, Inches(0.65), bar_t + Inches(0.15), Inches(3.5), Inches(0.7),
            '$180K–$450K', 36, bold=True, color=YELLOW)
add_textbox(s, Inches(0.65), bar_t + Inches(0.85), bar_w - Inches(0.3), Inches(0.75),
            'Annual platform licensing cost avoided — BambooHR/Rippling $6–$9/user/mo '
            '($108K–$162K/yr basic HR only). Workflow platforms $10–$25/user/mo = '
            '$180K–$450K/yr at 1,500 employees. This runs on Google Workspace TEAM Group already pays for.',
            10, color=TEXT)

# Right bar — outsource cost
add_rect(s, Inches(6.8), bar_t, bar_w, bar_h, fill_color=RGBColor(0x04, 0x18, 0x0c))
add_textbox(s, Inches(6.95), bar_t + Inches(0.15), Inches(3.5), Inches(0.7),
            '$65K–$115K', 36, bold=True, color=GREEN)
add_textbox(s, Inches(6.95), bar_t + Inches(0.85), bar_w - Inches(0.3), Inches(0.75),
            'External build cost avoided — outsourcing to a third-party developer would have '
            'cost $65,000–$115,000 as a one-time contract, before ongoing maintenance. '
            'Built internally using existing tools.',
            10, color=TEXT)

# Comparison table
add_table(s,
    l=Inches(0.5), t=Inches(3.7), w=Inches(12.33),
    headers=['Approach', 'Build / Setup Cost', 'Annual Recurring', 'Data Ownership', 'Workflow Tracking'],
    rows=[
        ['This Project',             '$0 external cost',   '$0',                         'Full — in your Drive', 'End-to-end, audited'],
        ['Outsource to contractor',  '$65K–$115K',         'Ongoing maintenance fees',   'Full',                 'Depends on scope'],
        ['SaaS Workflow Platform',   'Setup + integration','$180K–$450K+ (1,500 users)','Vendor-hosted',        'Yes, at significant cost'],
    ],
    hl_rows=[0],
    col_widths=[Inches(2.6), Inches(2.2), Inches(2.6), Inches(2.2), Inches(2.75)],
)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — Why This Is Noteworthy
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
set_bg(s, BG)
add_header(s, 'Why This Is Noteworthy')
add_slide_title(s, 'Why This Is Noteworthy',
                'Enterprise workflow capabilities — built on tools you already own.')

add_textbox(s, Inches(0.62), Inches(1.85), Inches(12.09), Inches(0.55),
            'Google Forms and Sheets are used across the organization for simple data collection. '
            'This system uses those same tools to deliver what typically requires a dedicated enterprise platform: '
            'a fully automated, multi-step approval workflow with enforcement, routing, and audit trail.',
            13, color=TEXT)

add_table(s,
    l=Inches(0.5), t=Inches(2.5), w=Inches(12.33),
    headers=['What most people use Google Forms for', 'What this system does'],
    rows=[
        ['Collect responses in a spreadsheet',
         'Trigger a multi-step approval chain — each step only unlocks after the prior one is complete, enforced automatically'],
        ['Send a single notification email',
         'Route the right email to the right person at each step; edits notify all affected parties with a full before/after change summary'],
        ['View raw submissions in a shared sheet',
         'Role-filtered live dashboard — HR sees all open requests, IT sees only their queue, Safety sees only what\'s relevant to them'],
        ['One form, one purpose',
         'Four complete workflow types, each with its own routing logic, task assignments, and approval gates — all managed in one system'],
    ],
    col_widths=[Inches(4.2), Inches(8.13)],
)

# Blue callout
add_rect(s, Inches(0.5), Inches(5.55), Inches(12.33), Inches(1.35), fill_color=RGBColor(0x04, 0x10, 0x20))
add_textbox(s, Inches(0.7), Inches(5.6), Inches(12.0), Inches(1.2),
            'No off-the-shelf product does this inside Google Workspace. Platforms with comparable multi-step '
            'approval routing — Kissflow ($10–$20/user/mo), Nintex ($25K–$60K+/yr), Process Street ($30/user/mo) — '
            'are external systems where your HR data lives on a vendor\'s servers. This delivers the same category '
            'of functionality at zero recurring cost, with data that never leaves Google Drive.',
            12, color=TEXT)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — Project Timeline (simplified Gantt as table)
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
set_bg(s, BG)
add_header(s, 'Project Timeline')
add_slide_title(s, 'Project Timeline',
                '5 months · solo · part-time — 331 tracked changes across 48 active days')

# Dev section header
add_textbox(s, Inches(0.62), Inches(1.88), Inches(3), Inches(0.28),
            'DEVELOPMENT', 10, bold=True, color=RGBColor(0xEB, 0x5C, 0x6D))

DEV_ROWS = [
    ['Foundation & New Hire',     'Dec 2025 – Jan 2026', 'Core pipeline · NEW_EMP forms · first prod launch'],
    ['Staging Environment',       'Feb 2026',            '3-environment separation · config architecture'],
    ['Termination / EOE',         'Mar 2026',            'Full TERM/EOE workflow · specialist routing'],
    ['Equipment Request',         'Mar – Apr 2026',      'Equipment + specialist task routing'],
    ['Status / Position Change',  'Apr 2026',            'Status change pipeline'],
    ['Action Items System',       'Apr 2026',            'Unified task model refactor'],
]

add_table(s,
    l=Inches(0.5), t=Inches(2.18), w=Inches(12.33),
    headers=['Stream', 'Period', 'Detail'],
    rows=DEV_ROWS,
    col_widths=[Inches(2.8), Inches(2.2), Inches(7.33)],
)

# QA section header
add_textbox(s, Inches(0.62), Inches(4.6), Inches(3), Inches(0.28),
            'QA & TESTING', 10, bold=True, color=BLUE)

QA_ROWS = [
    ['New Hire QA',         'Jan 23, 2026',  '~3.5 h',   '72 workflows submitted — first prod launch'],
    ['Staging Setup QA',    'Feb 24, 2026',  '~3 h',     '3 cross-environment workflow tests'],
    ['TERM/EOE Sprint',     'Mar 6–9, 2026', '~7.5 h',   '30 TERM/EOE workflows — gating & routing confirmed'],
    ['Full Regression',     'Apr 21–22',     '~14 h',    'All 4 workflow types · Action Items · Equipment · Status Change'],
    ['Pre-launch Verify',   'Apr 27–28',     '~12.5 h',  '38 dev deploys · all gates signed off'],
]

add_table(s,
    l=Inches(0.5), t=Inches(4.9), w=Inches(12.33),
    headers=['Session', 'Date', 'Hours', 'Scope / Output'],
    rows=QA_ROWS,
    col_widths=[Inches(2.2), Inches(1.8), Inches(1.0), Inches(7.33)],
)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — Development & Testing (phase table)
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
set_bg(s, BG)
add_header(s, 'Development & Testing')
add_slide_title(s, 'Development & Testing',
                '~125 h development · ~40 h QA · 105 staging workflows submitted')

add_table(s,
    l=Inches(0.5), t=Inches(1.85), w=Inches(12.33),
    headers=['Phase', 'Period', 'Tracked Changes', 'Dev Hours', 'QA Hours', 'Total Hours', 'QA Submissions / Output'],
    rows=[
        ['Foundation & Architecture',                     'Dec 2025',       '20',   '~15 h',  '—',       '~15 h',   '—'],
        ['New Hire Workflow',                              'Jan 2026',       '23',   '~18 h',  '~3.5 h',  '~22 h',   '72 NEW_EMP workflows — first prod launch'],
        ['Staging Environment',                           'Feb 2026',       '17',   '~6.5 h', '~3 h',    '~9.5 h',  '3 cross-environment tests'],
        ['Termination / EOE + Equipment',                 'Mar 2026',       '94',   '~36 h',  '~7.5 h',  '~44 h',   '30 TERM/EOE workflows submitted'],
        ['Action Items · Status Change · Equip QA · Pre-launch', 'Apr 2026','177',  '~49 h',  '~26 h',   '~75 h',   'Full regression all 4 types · 38 deploys'],
        ['Total',                                         'Dec 2025–Apr 2026','331','~125 h', '~40 h',   '~174 h',  '105 staging workflows · ~305 sheet writes · 4 types live'],
    ],
    hl_rows=[3, 4],
    total_row=True,
    col_widths=[Inches(2.7), Inches(1.5), Inches(1.3), Inches(1.1), Inches(1.1), Inches(1.1), Inches(3.52)],
)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — Maintainability & Risk
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
set_bg(s, BG)
add_header(s, 'Maintainability & Risk')
add_slide_title(s, 'Maintainability & Risk',
                'Version-controlled, environment-isolated, safe to hand off.')

add_table(s,
    l=Inches(0.5), t=Inches(1.85), w=Inches(12.33),
    headers=['Concern', 'How It\'s Addressed'],
    rows=[
        ['Single developer dependency',
         'All logic is version-controlled and documented. Separate testing environments make it safe '
         'for any developer to extend without risk to live data or active workflows.'],
        ['Updating the system',
         'Changes are tested in a staging environment first and only promoted to production after validation. '
         'Data model updates run safely without taking the system offline.'],
        ['Google Workspace dependency',
         'No new external dependency — the risk profile is identical to any existing use of Google Forms '
         'or Sheets in the organization.'],
        ['AI-generated code quality',
         'Every component was reviewed, tested against real HR scenarios, and validated before going live. '
         'The ~174 hours is primarily review, direction, and QA — not passive acceptance of AI output.'],
    ],
    col_widths=[Inches(3.0), Inches(9.33)],
)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — Bottom Line
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
set_bg(s, BG)
add_header(s, 'Bottom Line')

# Centred heading
add_textbox(s, Inches(1.0), Inches(1.2), Inches(11.33), Inches(0.7),
            'Bottom Line', 40, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# Red underline accent
add_rect(s, Inches(5.67), Inches(2.0), Inches(2.0), Inches(0.05), fill_color=RED)

# Green callout box
add_rect(s, Inches(1.0), Inches(2.2), Inches(11.33), Inches(2.5), fill_color=RGBColor(0x04, 0x18, 0x0c))
add_textbox(s, Inches(1.2), Inches(2.3), Inches(11.0), Inches(2.3),
            "TEAM Group's HR onboarding and offboarding process moved from emailed spreadsheets to a "
            "fully tracked, accountable, automated workflow — built on tools the organization already uses, "
            "with zero new software, zero recurring licensing cost, and no external contractors.\n\n"
            "The equivalent capability would cost $65K–$115K to outsource to a contractor and "
            "$180K–$450K+ per year to license from a workflow automation platform.",
            16, color=TEXT)

# Stat pills
PILLS = ['331 tracked changes', '48 active days', '~174 hours total', '4 workflow types live']
pill_w = Inches(2.7)
gap_p  = Inches(0.15)
total_pw = len(PILLS) * pill_w + (len(PILLS) - 1) * gap_p
start_l = (W - total_pw) / 2

for i, pill in enumerate(PILLS):
    l = start_l + i * (pill_w + gap_p)
    add_rect(s, l, Inches(5.0), pill_w, Inches(0.5), fill_color=RGBColor(0x1e, 0x1e, 0x1e))
    add_textbox(s, l, Inches(5.02), pill_w, Inches(0.46),
                pill, 13, bold=True, color=MUTED, align=PP_ALIGN.CENTER)


# ── Save ─────────────────────────────────────────────────────────────────────
out_path = r"P:\Repos\github\danieljbinns\APP SCRIPT FORMS\employee_forms_deployment\PROJECT_SUMMARY_EXEC_SLIDES.pptx"
prs.save(out_path)
print(f"Saved -> {out_path}")
print(f"Slides: {len(prs.slides)}")
